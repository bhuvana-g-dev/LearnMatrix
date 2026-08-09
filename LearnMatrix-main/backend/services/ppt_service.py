"""
services/ppt_service.py

Builds a downloadable .pptx "Study Summary" deck. Three modes,
mirroring Flashcards' mode selector:
  - generate_study_summary_pptx: from an already-generated
    learning_notes entry (services/notes_repository.py)
  - generate_chat_summary_pptx: from one saved AI Chat session's
    Q&A turns (services/chat_repository.py)
  - generate_sources_summary_pptx: from the user's uploaded/linked
    chat sources (services/embedding_service.py)

No LLM call in ANY of the three modes — each just reshapes data that
already exists into a {title, summary, sections, keyTakeaways} shape,
which _build_pptx_from_notes turns into slides. Regenerating that text
through an agent would be a slower, costlier way to get content already
on hand.

Uses python-pptx to build the deck in memory (BytesIO) — nothing is
written to disk, so there's no cleanup/temp-file lifecycle to manage;
routes/ppt_routes.py streams the bytes straight back as a file download.
"""

from io import BytesIO

from pptx import Presentation
from pptx.util import Pt

from firebase.firebase_config import get_firestore_client
from services import chat_repository, embedding_service, notes_repository

TITLE_AND_CONTENT_LAYOUT = 1


class PptServiceError(Exception):
    pass


def generate_study_summary_pptx(skill: str, topic: str, focus_band: str) -> tuple[BytesIO, str]:
    """Raises PptServiceError if no notes exist yet for this
    (skill, topic, focus_band) — same precondition as
    embedding_service.index_learning_notes."""
    db = get_firestore_client()
    notes = notes_repository.get_cached_notes(db, skill, topic, focus_band)
    if not notes:
        raise PptServiceError(
            f"No study notes found yet for '{skill} / {topic}' ({focus_band}). "
            "Open that topic in the Learning Hub first so notes are generated."
        )
    safe_topic = _safe_filename(topic)
    return _build_pptx_from_notes(notes, subtitle=f"{focus_band.title()} level"), f"{safe_topic}_study_summary.pptx"


def generate_sources_summary_pptx(uid: str) -> tuple[BytesIO, str]:
    db = get_firestore_client()
    sources_content = embedding_service.get_sources_with_text(db, uid)
    if not sources_content:
        raise PptServiceError("No sources found yet — upload a source first.")

    notes = {
        "title": "Your Sources",
        "summary": "",
        "sections": [{"heading": s["title"], "content": s["text"]} for s in sources_content],
        "keyTakeaways": [],
    }
    return _build_pptx_from_notes(notes, subtitle="Study Summary from Sources"), "sources_study_summary.pptx"


def generate_chat_summary_pptx(uid: str, session_id: str) -> tuple[BytesIO, str]:
    if not session_id:
        raise PptServiceError("No chat conversation selected — open or start a chat first.")
    db = get_firestore_client()
    history = chat_repository.get_session_messages(db, uid, session_id, limit=0)
    if not history:
        raise PptServiceError("That conversation has no messages yet — ask the AI Study Assistant something first.")

    sections = []
    for i, turn in enumerate(history):
        if turn.get("role") == "user":
            answer = history[i + 1].get("content", "") if i + 1 < len(history) and history[i + 1].get("role") == "assistant" else ""
            sections.append({"heading": turn.get("content", "")[:60], "content": answer})

    notes = {"title": "Your Chat", "summary": "", "sections": sections, "keyTakeaways": []}
    return _build_pptx_from_notes(notes, subtitle="Study Summary from AI Chat"), "chat_study_summary.pptx"


def generate_custom_text_pptx(text: str) -> tuple[BytesIO, str]:
    if not text or not text.strip():
        raise PptServiceError("Type something first.")
    title = text.strip().split("\n")[0][:60] or "Your Topic"
    notes = {"title": title, "summary": "", "sections": [{"heading": "Your Input", "content": text.strip()}], "keyTakeaways": []}
    return _build_pptx_from_notes(notes, subtitle="Study Summary"), "custom_study_summary.pptx"


def _build_pptx_from_notes(notes: dict, subtitle: str) -> BytesIO:
    prs = Presentation()
    _add_title_slide(prs, notes.get("title") or "Study Summary", subtitle)

    if notes.get("summary"):
        _add_text_slide(prs, "Summary", notes["summary"])

    for section in notes.get("sections", []):
        heading = section.get("heading", "Section")
        content = section.get("content", "")
        if content:
            _add_text_slide(prs, heading, content)

    takeaways = notes.get("keyTakeaways", [])
    if takeaways:
        _add_bullet_slide(prs, "Key Takeaways", takeaways)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def _safe_filename(value: str) -> str:
    return "".join(c if c.isalnum() or c in " -_" else "_" for c in value)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"LearnMatrix \u2014 {subtitle}"


def _add_text_slide(prs: Presentation, heading: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = heading
    body_placeholder = slide.placeholders[1]
    text_frame = body_placeholder.text_frame
    text_frame.word_wrap = True
    # A single dense paragraph reads poorly on a slide, so split on
    # sentence boundaries into shorter paragraphs rather than dumping
    # the whole block into one text run.
    sentences = [s.strip() for s in body.replace("\n", " ").split(". ") if s.strip()]
    if not sentences:
        text_frame.text = body[:2000]
        return
    text_frame.text = sentences[0] + ("." if not sentences[0].endswith(".") else "")
    for sentence in sentences[1:20]:  # cap paragraphs per slide — very long sources shouldn't produce one giant slide
        p = text_frame.add_paragraph()
        p.text = sentence + ("." if not sentence.endswith(".") else "")
        p.font.size = Pt(18)


def _add_bullet_slide(prs: Presentation, heading: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = heading
    text_frame = slide.placeholders[1].text_frame
    text_frame.word_wrap = True
    text_frame.text = f"\u2022 {bullets[0]}"
    for b in bullets[1:]:
        p = text_frame.add_paragraph()
        p.text = f"\u2022 {b}"
        p.font.size = Pt(18)
