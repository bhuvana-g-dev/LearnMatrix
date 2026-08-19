"""
services/ppt_service.py

Builds a downloadable .pptx "Study Summary" deck, styled with the
LearnMatrix navy/gold brand palette (colored backgrounds, accent bars,
badge-numbered sections, colored bullet markers) instead of the
default plain PowerPoint template. Three modes, mirroring Flashcards'
mode selector:
  - generate_study_summary_pptx: from an already-generated
    learning_notes entry (services/notes_repository.py)
  - generate_chat_summary_pptx: from one saved AI Chat session's
    Q&A turns (services/chat_repository.py)
  - generate_sources_summary_pptx: from the user's uploaded/linked
    chat sources (services/embedding_service.py)

No LLM call in ANY of the three modes — each just reshapes data that
already exists into a {title, summary, sections, keyTakeaways} shape,
which _build_pptx_from_notes turns into slides. Regenerating that text
through an agent would be a slower, costlier way to get content already
on hand.

There's no external image-search/stock-photo API configured in this
backend (no Unsplash/Pexels key in config/settings.py), so "visual"
content here means brand-colored shapes/badges/accent bars drawn
directly with python-pptx — not photographs. Wiring in a real photo
API is a separate, later change (would need an API key added to
Settings).

Uses python-pptx to build the deck in memory (BytesIO) — nothing is
written to disk, so there's no cleanup/temp-file lifecycle to manage;
routes/ppt_routes.py streams the bytes straight back as a file download.
"""

from io import BytesIO

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from firebase.firebase_config import get_firestore_client
from services import chat_repository, embedding_service, notes_repository

# --- LearnMatrix brand palette (mirrors frontend/src/constants/theme.js) ---
NAVY = RGBColor(0x0D, 0x1B, 0x3D)
NAVY_MID = RGBColor(0x3E, 0x4A, 0x66)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GOLD_LIGHT = RGBColor(0xE8, 0xB9, 0x3D)
CREAM = RGBColor(0xFB, 0xF3, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Alternating accent color per content slide, cycling through the palette
# so a long deck doesn't look monotone.
ACCENT_CYCLE = [GOLD, NAVY, GOLD_LIGHT]


class PptServiceError(Exception):
    pass


def generate_study_summary_pptx(skill: str, topic: str, focus_band: str) -> tuple[BytesIO, str]:
    """Raises PptServiceError if no notes exist yet for this
    (skill, topic, focus_band) — same precondition as
    embedding_service.index_learning_notes."""
    db = get_firestore_client()
    notes = notes_repository.get_cached_notes(db, skill, topic, focus_band)
    if not notes:
        raise PptServiceError(
            f"No study notes found yet for '{skill} / {topic}' ({focus_band}). "
            "Open that topic in the Learning Hub first so notes are generated."
        )
    safe_topic = _safe_filename(topic)
    enriched = _enrich_notes_for_deck(notes, label=f"{skill} - {topic}") or notes
    return _build_pptx_from_notes(enriched, subtitle=f"{focus_band.title()} level"), f"{safe_topic}_study_summary.pptx"


def generate_sources_summary_pptx(uid: str) -> tuple[BytesIO, str]:
    db = get_firestore_client()
    sources_content = embedding_service.get_sources_with_text(db, uid)
    if not sources_content:
        raise PptServiceError("No sources found yet — upload a source first.")

    notes = {
        "title": "Your Sources",
        "summary": "",
        "sections": [{"heading": s["title"], "content": s["text"]} for s in sources_content],
        "keyTakeaways": [],
    }
    enriched = _enrich_notes_for_deck(notes, label="these sources") or notes
    return _build_pptx_from_notes(enriched, subtitle="Study Summary from Sources"), "sources_study_summary.pptx"


def generate_chat_summary_pptx(uid: str, session_id: str) -> tuple[BytesIO, str]:
    if not session_id:
        raise PptServiceError("No chat conversation selected — open or start a chat first.")
    db = get_firestore_client()
    history = chat_repository.get_session_messages(db, uid, session_id, limit=0)
    if not history:
        raise PptServiceError("That conversation has no messages yet — ask the AI Study Assistant something first.")

    sections = []
    for i, turn in enumerate(history):
        if turn.get("role") == "user":
            answer = history[i + 1].get("content", "") if i + 1 < len(history) and history[i + 1].get("role") == "assistant" else ""
            sections.append({"heading": turn.get("content", "")[:60], "content": answer})

    notes = {"title": "Your Chat", "summary": "", "sections": sections, "keyTakeaways": []}
    enriched = _enrich_notes_for_deck(notes, label="this conversation") or notes
    return _build_pptx_from_notes(enriched, subtitle="Study Summary from AI Chat"), "chat_study_summary.pptx"


def generate_custom_text_pptx(text: str) -> tuple[BytesIO, str]:
    """AI-expands the student's short prompt/notes into a full deck
    (see services/slide_deck_service.py) rather than dumping the raw
    text onto a single slide — same idea as Gamma/NotebookLM turning a
    one-line prompt into a real presentation."""
    if not text or not text.strip():
        raise PptServiceError("Type something first.")
    from services.slide_deck_service import generate_deck_content, SlideDeckServiceError

    try:
        notes = generate_deck_content(text.strip())
    except SlideDeckServiceError as exc:
        raise PptServiceError(str(exc)) from exc
    return build_pptx_from_deck_content(notes, subtitle="Study Summary"), f"{_safe_filename(notes.get('title') or 'custom')}_study_summary.pptx"


def _enrich_notes_for_deck(notes: dict, label: str) -> dict | None:
    """Runs the SAME Gamma/NotebookLM-style AI enrichment the "type a
    topic" custom deck already gets (agents/slide_deck_agent.py, via
    services/slide_deck_service.py) on top of already-existing plain
    notes/chat/sources content — so every export mode gets varied
    layouts (list/process/comparison), per-item icons, and real
    AI-generated images, not just the custom-prompt one.

    Returns None (caller falls back to the original plain notes) if the
    AI call fails for any reason — a slow/unavailable LLM should
    degrade an export to "plain but still downloadable", never break it
    outright. This does mean every pptx download now makes an LLM call
    plus per-section image calls — see services/slide_deck_service.py's
    _attach_section_images, which is sequential per section, so a
    gunicorn worker timeout that's too short (see routes/ai_assessment_routes.py
    and routes/slidedeck_routes.py's own timeout notes) will surface here too."""
    from services.slide_deck_service import generate_deck_content, SlideDeckServiceError

    combined = (notes.get("title") or "").strip() + "\n\n"
    for s in notes.get("sections", []):
        combined += f"## {s.get('heading', '')}\n{s.get('content', '')}\n\n"
    combined = combined.strip()
    if not combined:
        return None

    try:
        return generate_deck_content(combined[:6000], label=label)
    except SlideDeckServiceError:
        return None


# ---------------------------------------------------------------------------
# Deck assembly — collects the notes dict content shared by every mode above
# into a `_DeckContent`-shaped list of sections, so the PDF service
# (services/pdf_service.py) can build a matching-design PDF from the exact
# same content without duplicating the Firestore-fetching logic in each
# generate_* function above.
# ---------------------------------------------------------------------------

def build_deck_sections(notes: dict) -> list[dict]:
    """Turns a {title, summary, sections, keyTakeaways} notes dict into
    the flat list of slide dicts both the pptx and pdf builders render,
    so the two file formats never drift apart. Each slide's "kind" is
    one of "text", "list", "comparison", "process", or "bullets" (Key
    Takeaways) — this is what gives the deck varied layouts instead of
    every slide being a plain heading+paragraph (see
    agents/slide_deck_agent.py, which tags each section's layout and
    per-item icon based on its content)."""
    slides = []
    if notes.get("summary"):
        slides.append({"kind": "text", "heading": "Summary", "body": notes["summary"]})

    for section in notes.get("sections", []):
        heading = section.get("heading", "Section")
        layout = section.get("layout", "text")

        if layout == "list" and isinstance(section.get("items"), list) and section["items"]:
            slides.append({"kind": "list", "heading": heading, "items": section["items"]})
        elif layout == "process" and isinstance(section.get("steps"), list) and section["steps"]:
            slides.append({"kind": "process", "heading": heading, "steps": section["steps"]})
        elif layout == "comparison" and isinstance(section.get("left"), dict) and isinstance(section.get("right"), dict):
            slides.append({"kind": "comparison", "heading": heading, "left": section["left"], "right": section["right"]})
        elif section.get("content"):
            slide = {"kind": "text", "heading": heading, "body": section["content"]}
            if section.get("image_url"):
                slide["image_url"] = section["image_url"]
            if section.get("subpoints"):
                slide["subpoints"] = section["subpoints"]
            slides.append(slide)

    takeaways = notes.get("keyTakeaways", [])
    if takeaways:
        slides.append({"kind": "bullets", "heading": "Key Takeaways", "items": takeaways})
    return slides


def _build_pptx_from_notes(notes: dict, subtitle: str) -> BytesIO:
    return build_pptx_from_deck_content(notes, subtitle)


def build_pptx_from_deck_content(notes: dict, subtitle: str) -> BytesIO:
    """Public entry point used both by the generate_*_pptx functions
    above (which fetch/produce `notes` themselves) and by
    routes/ppt_routes.py's "from-content" endpoint, which receives an
    already-generated deck (e.g. from the AI slide-deck preview the
    student already looked at) and just needs it rendered — no second
    LLM call, so the downloaded file always matches what was previewed."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # fully blank — every element below is hand-placed

    _add_title_slide(prs, blank_layout, notes.get("title") or "Study Summary", subtitle)

    deck_sections = build_deck_sections(notes)
    for i, slide_data in enumerate(deck_sections):
        accent = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        kind = slide_data["kind"]
        if kind == "bullets":
            _add_bullet_slide(prs, blank_layout, slide_data["heading"], slide_data["items"], i + 1, accent)
        elif kind == "list":
            _add_list_slide(prs, blank_layout, slide_data["heading"], slide_data["items"], i + 1, accent)
        elif kind == "process":
            _add_process_slide(prs, blank_layout, slide_data["heading"], slide_data["steps"], i + 1, accent)
        elif kind == "comparison":
            _add_comparison_slide(prs, blank_layout, slide_data["heading"], slide_data["left"], slide_data["right"], i + 1)
        else:
            _add_text_slide(prs, blank_layout, slide_data["heading"], slide_data["body"], i + 1, accent, slide_data.get("image_url"), slide_data.get("subpoints"))

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def _safe_filename(value: str) -> str:
    return "".join(c if c.isalnum() or c in " -_" else "_" for c in value)


# ---------------------------------------------------------------------------
# Slide builders — every shape is hand-placed on a blank layout so the
# design (colored backgrounds, accent bars, numbered badges, bullet
# markers) is fully controlled rather than inherited from a default
# PowerPoint theme.
# ---------------------------------------------------------------------------

def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _distribute_fill(n: int, avail_h_in: float, gap_in: float, min_item: float, max_item: float) -> tuple[float, float]:
    """Shared fix for the "few items -> big blank strip at the bottom
    of the slide" problem across list/bullet/comparison layouts. Given
    n stacked items and the available vertical space, grows each
    item's height to fill that space (clamped to [min_item, max_item]
    so a 1-2 item slide doesn't get one absurdly tall card), and
    returns a top offset so any leftover space is centered rather than
    left as dead space below the last item. Returns (item_height_in,
    top_offset_in), both in inches."""
    if n <= 0:
        return min_item, 0.0
    ideal = (avail_h_in - gap_in * (n - 1)) / n
    item_h_in = max(min_item, min(ideal, max_item))
    content_h_in = item_h_in * n + gap_in * (n - 1)
    offset_in = max(0.0, (avail_h_in - content_h_in) / 2)
    return item_h_in, offset_in


def _add_title_slide(prs: Presentation, layout, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, NAVY)
    bg.shadow.inherit = False

    # Large soft gold arc/circle bleeding off the right edge as a decorative
    # accent — stands in for photography since no image-search API is
    # configured (see module docstring).
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-2.5), Inches(7), Inches(7))
    _set_fill(circle, GOLD)
    circle.fill.fore_color.brightness = 0.0
    circle.shadow.inherit = False
    _set_transparency(circle, 70)

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(3.6), Inches(3.4), Inches(3.4))
    _set_fill(circle2, GOLD_LIGHT)
    circle2.shadow.inherit = False
    _set_transparency(circle2, 40)

    # Thin gold rule above the title, LearnMatrix wordmark styling.
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.55), Inches(1.1), Pt(4))
    _set_fill(rule, GOLD)
    rule.shadow.inherit = False

    brand_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(6), Inches(0.5))
    brand_tf = brand_box.text_frame
    brand_tf.text = "LEARNMATRIX"
    brand_run = brand_tf.paragraphs[0].runs[0]
    brand_run.font.size = Pt(15)
    brand_run.font.bold = True
    brand_run.font.color.rgb = GOLD_LIGHT
    brand_run.font.name = "Arial"

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(2.9), Inches(9.5), Inches(2.2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.text = title
    title_run = title_tf.paragraphs[0].runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = WHITE
    title_run.font.name = "Arial"

    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(9), Inches(0.6))
    sub_tf = sub_box.text_frame
    sub_tf.text = f"{subtitle}"
    sub_run = sub_tf.paragraphs[0].runs[0]
    sub_run.font.size = Pt(18)
    sub_run.font.color.rgb = CREAM
    sub_run.font.name = "Arial"


def _add_slide_chrome(prs, slide, heading, index, accent):
    """Shared header used by every content slide: colored left accent bar,
    a numbered badge, and the heading — keeps _add_text_slide and
    _add_bullet_slide visually consistent."""
    # Full-height accent bar down the left edge.
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), prs.slide_height)
    _set_fill(bar, accent)
    bar.shadow.inherit = False

    # White page background.
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), 0, prs.slide_width - Inches(0.35), prs.slide_height)
    _set_fill(bg, WHITE)
    bg.shadow.inherit = False
    # send background behind everything added after it by re-adding chrome on top — python-pptx has no z-order API,
    # so shapes are simply added in back-to-front order (bar/bg first, since they're added first, is already correct).

    # Numbered badge.
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(0.55), Inches(0.55), Inches(0.55))
    _set_fill(badge, accent)
    badge.shadow.inherit = False
    badge_tf = badge.text_frame
    badge_tf.word_wrap = False
    badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_tf.text = str(index)
    badge_run = badge_tf.paragraphs[0].runs[0]
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge_run.font.size = Pt(18)
    badge_run.font.bold = True
    badge_run.font.color.rgb = WHITE if accent != GOLD_LIGHT else NAVY
    badge_run.font.name = "Arial"

    heading_box = slide.shapes.add_textbox(Inches(1.55), Inches(0.5), Inches(11), Inches(0.75))
    heading_tf = heading_box.text_frame
    heading_tf.word_wrap = True
    heading_tf.text = heading
    heading_run = heading_tf.paragraphs[0].runs[0]
    heading_run.font.size = Pt(26)
    heading_run.font.bold = True
    heading_run.font.color.rgb = NAVY
    heading_run.font.name = "Arial"

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.55), Inches(1.35), Inches(1.4), Pt(3))
    _set_fill(rule, accent)
    rule.shadow.inherit = False


def _add_text_slide(prs: Presentation, layout, heading: str, body: str, index: int, accent, image_url: str | None = None, subpoints: list | None = None) -> None:
    slide = prs.slides.add_slide(layout)
    _add_slide_chrome(prs, slide, heading, index, accent)

    # When a photo/illustration is available for this section (AI-
    # generated or Pexels — see services/slide_deck_service.py), the
    # text column narrows to make room for it on the right instead of
    # running full-width. Sub-point highlights (see
    # agents/slide_deck_agent.py) share that same right-hand column,
    # stacking below the image when both are present.
    image_bytes = None
    if image_url:
        from services.image_service import fetch_image_bytes
        image_bytes = fetch_image_bytes(image_url)

    has_sidebar = bool(image_bytes) or bool(subpoints)
    text_width = Inches(6.6) if has_sidebar else Inches(11.1)
    content_h_in = 5.35

    sentences = [s.strip() for s in body.replace("\n", " ").split(". ") if s.strip()]
    if not sentences:
        sentences = [body[:2000]]

    # No sidebar (no image/subpoints) means the paragraph is the ONLY
    # thing on the slide, so a short section (a demoted-to-text
    # section can legitimately be just 1-2 sentences — see
    # agents/slide_deck_agent.py) used to leave a big blank strip
    # below it. Scale font size UP and center the block vertically so
    # a short paragraph still reads as an intentional, full slide
    # instead of a stub sitting in the top-left corner.
    if has_sidebar:
        body_font_size = Pt(16) if len(sentences) > 5 else Pt(18)
        body_top_in = 1.7
    else:
        if len(sentences) <= 3:
            body_font_size = Pt(30)
        elif len(sentences) <= 5:
            body_font_size = Pt(22)
        else:
            body_font_size = Pt(18)
        # Rough line-count estimate to vertically center the block —
        # doesn't need to be exact, just close enough that a short
        # paragraph isn't glued to the top of an otherwise empty slide.
        chars_per_line = max(20, int(text_width / 914400 * 96 / (body_font_size.pt * 0.55)))
        total_chars = sum(len(s) for s in sentences)
        est_lines = max(len(sentences), -(-total_chars // chars_per_line))
        est_h_in = est_lines * (body_font_size.pt / 72) * 1.7
        body_top_in = 1.7 + max(0, (content_h_in - est_h_in) / 2)

    body_box = slide.shapes.add_textbox(Inches(1.55), Inches(body_top_in), text_width, Inches(content_h_in))
    text_frame = body_box.text_frame
    text_frame.word_wrap = True

    # A single dense paragraph reads poorly on a slide, so split on
    # sentence boundaries into shorter paragraphs rather than dumping
    # the whole block into one text run.
    first = True
    for sentence in sentences[:26]:  # cap paragraphs per slide — very long sources shouldn't produce one giant slide
        text = sentence + ("." if not sentence.endswith(".") else "")
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = text
        run.font.size = body_font_size
        run.font.color.rgb = NAVY_MID
        run.font.name = "Arial"

    sidebar_top = Inches(1.85)
    if image_bytes:
        img_left, img_top, img_w, img_h = Inches(8.45), sidebar_top, Inches(4.2), Inches(3.15)
        # A plain add_picture() would be a hard-edged rectangle sitting on
        # a page of otherwise all-rounded shapes — a soft accent frame
        # behind it (peeking out on two sides) reads as an intentional
        # framed photo instead of a raw pasted image.
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_left - Inches(0.12), img_top - Inches(0.12), img_w + Inches(0.24), img_h + Inches(0.24))
        frame.adjustments[0] = 0.04
        _set_fill(frame, accent)
        frame.shadow.inherit = False
        try:
            slide.shapes.add_picture(BytesIO(image_bytes), img_left, img_top, width=img_w, height=img_h)
            sidebar_top = img_top + img_h + Inches(0.4)
        except Exception:
            pass  # a corrupt/unsupported download shouldn't break the whole deck — the accent frame alone still looks intentional

    if subpoints:
        _add_key_points_panel(slide, subpoints, Inches(8.45), sidebar_top, Inches(4.2), accent)


def _add_key_points_panel(slide, subpoints: list, x, y, width, accent) -> None:
    """A compact 'Key Points' card stacked in the text slide's sidebar
    (below the image when one is present) — short highlight bullets
    (see agents/slide_deck_agent.py's "subpoints") give a Gamma/
    NotebookLM-style at-a-glance summary alongside the long-form
    paragraph, instead of forcing the reader through the whole
    paragraph to find the key facts."""
    label_box = slide.shapes.add_textbox(x, y, width, Inches(0.35))
    lf = label_box.text_frame
    lf.word_wrap = True
    lp = lf.paragraphs[0]
    lrun = lp.add_run()
    lrun.text = "KEY POINTS"
    lrun.font.size = Pt(12)
    lrun.font.bold = True
    lrun.font.name = "Arial"
    lrun.font.color.rgb = accent

    row_y = y + Inches(0.38)
    for sp in subpoints[:5]:
        text, _icon = _item_text_icon(sp)
        if not text:
            continue
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, row_y + Inches(0.06), Inches(0.09), Inches(0.09))
        _set_fill(dot, accent)
        dot.shadow.inherit = False
        dot.line.fill.background()

        item_box = slide.shapes.add_textbox(x + Inches(0.22), row_y - Inches(0.05), width - Inches(0.22), Inches(0.6))
        tf = item_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(13)
        run.font.name = "Arial"
        run.font.color.rgb = NAVY_MID

        line_count = max(1, -(-len(text) // 42))  # rough wrap estimate, matches ~4.2" width at 13pt
        row_y += Inches(0.24 * line_count + 0.16)


# icon tag (see agents/slide_deck_agent.py's ICON_VOCAB) -> a built-in
# python-pptx autoshape that reads as that icon at a glance, so list
# items and key-takeaway cards vary by MEANING instead of every card
# using the same numbered circle.
ICON_SHAPE_MAP = {
    "check": MSO_SHAPE.OVAL,
    "star": MSO_SHAPE.STAR_5_POINT,
    "warning": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "gear": MSO_SHAPE.GEAR_6,
    "database": MSO_SHAPE.CAN,
    "network": MSO_SHAPE.HEXAGON,
    "shield": MSO_SHAPE.PENTAGON,
    "zap": MSO_SHAPE.LIGHTNING_BOLT,
    "cloud": MSO_SHAPE.CLOUD,
    "book": MSO_SHAPE.OVAL,
}


def _icon_badge(slide, icon: str, x, y, size, fill_color):
    """Adds one icon-shaped badge (see ICON_SHAPE_MAP) filled with
    fill_color. "check" gets a checkmark glyph overlaid — every other
    icon's shape alone (star, triangle, gear, hexagon, ...) is
    recognizable without needing a text glyph on top, which sidesteps
    relying on emoji/symbol font support in PowerPoint/LibreOffice."""
    shape_type = ICON_SHAPE_MAP.get(icon, MSO_SHAPE.OVAL)
    shape = slide.shapes.add_shape(shape_type, x, y, size, size)
    _set_fill(shape, fill_color)
    shape.shadow.inherit = False
    if icon in ("check", "book"):
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = "\u2713" if icon == "check" else ""
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if tf.paragraphs[0].runs:
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(max(10, int(size.pt * 0.42)))
            run.font.bold = True
            run.font.color.rgb = WHITE if fill_color != GOLD_LIGHT else NAVY
    return shape


def _item_text_icon(item) -> tuple[str, str]:
    """List items / key takeaways are {"text","icon"} dicts once past
    the agent's normalization, but this also accepts a plain string
    (e.g. an older cached deck) so nothing breaks on legacy content."""
    if isinstance(item, dict):
        return item.get("text", ""), item.get("icon", "check")
    return str(item), "check"


def _add_bullet_slide(prs: Presentation, layout, heading: str, bullets: list, index: int, accent) -> None:
    """Key Takeaways — each point in its own soft-tinted rounded card
    with a checkmark badge, instead of a plain circle+text bullet line."""
    slide = prs.slides.add_slide(layout)
    _add_slide_chrome(prs, slide, heading, index, accent)

    # Card height scales to fill down to the bottom margin (see
    # _distribute_fill) — a deck with only 4-5 takeaways used to leave
    # a tall blank gap under the last card.
    top_in, gap_in = 1.85, 0.18
    card_h_in, offset_in = _distribute_fill(len(bullets), 7.5 - top_in - 0.5, gap_in, min_item=0.72, max_item=1.35)
    top = Inches(top_in + offset_in)
    card_h = Inches(card_h_in)
    icon_size = Inches(0.4) if card_h_in <= 0.9 else Inches(0.5)
    for b in bullets:
        text, icon = _item_text_icon(b)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.55), top, Inches(10.9), card_h)
        card.adjustments[0] = 0.18
        _set_fill(card, CREAM)
        card.shadow.inherit = False

        _icon_badge(slide, icon, Inches(1.75), top + (card_h - icon_size) / 2, icon_size, accent)

        text_box = slide.shapes.add_textbox(Inches(2.35), top + Inches(0.06), Inches(9.9), card_h - Inches(0.1))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = text
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(16) if card_h_in <= 0.9 else Pt(18)
        run.font.bold = True
        run.font.color.rgb = NAVY
        run.font.name = "Arial"
        top += card_h + Inches(gap_in)


def _add_list_slide(prs: Presentation, layout, heading: str, items: list, index: int, accent) -> None:
    """Features/steps/examples — a grid of colored icon cards instead
    of a plain bulleted paragraph, e.g. for a "Key Features" section.
    Each item's icon (see agents/slide_deck_agent.py) is picked per its
    own meaning rather than a generic sequence number.

    Card height/row-gap SCALE to fill the available vertical space down
    to the bottom margin (see _distribute_fill) instead of a fixed
    1.05in — a 3-4 item list used to leave a big blank strip at the
    bottom of the slide; now the grid grows to occupy the full content
    area (with any small leftover centered) regardless of item count."""
    slide = prs.slides.add_slide(layout)
    _add_slide_chrome(prs, slide, heading, index, accent)

    cols = 2 if len(items) <= 4 else 3
    gap_in = 0.3
    area_left, area_top_in = Inches(1.55), 1.85
    area_w = prs.slide_width - area_left - Inches(0.6)
    card_w = (area_w - Inches(gap_in) * (cols - 1)) / cols
    rows = -(-len(items) // cols)  # ceil
    avail_h_in = 7.5 - area_top_in - 0.5  # down to bottom margin
    card_h_in, offset_in = _distribute_fill(rows, avail_h_in, gap_in, min_item=1.05, max_item=2.1)
    card_h = Inches(card_h_in)
    area_top = Inches(area_top_in + offset_in)

    for i, item in enumerate(items):
        text, icon = _item_text_icon(item)
        r, c = divmod(i, cols)
        x = area_left + c * (card_w + Inches(gap_in))
        y = area_top + r * (card_h + Inches(gap_in))

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.adjustments[0] = 0.12
        _set_fill(card, CREAM)
        card.line.color.rgb = accent
        card.line.width = Pt(1)
        card.shadow.inherit = False

        _icon_badge(slide, icon, x + Inches(0.18), y + Inches(0.18), Inches(0.4), accent)

        text_box = slide.shapes.add_textbox(x + Inches(0.72), y + Inches(0.1), card_w - Inches(0.88), card_h - Inches(0.2))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = text
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = NAVY
        run.font.name = "Arial"


def _add_process_slide(prs: Presentation, layout, heading: str, steps: list, index: int, accent) -> None:
    """"How it works" / ordered setup steps — a left-to-right chip flow
    with numbered circles and arrow connectors, instead of a bulleted
    list that doesn't read as a sequence."""
    slide = prs.slides.add_slide(layout)
    _add_slide_chrome(prs, slide, heading, index, accent)

    n = len(steps)
    area_left = Inches(1.55)
    area_w = prs.slide_width - area_left - Inches(0.6)
    arrow_w = Inches(0.45)
    # Chip height GROWS for a short sequence (2-3 steps), and the whole
    # row is centered in the available vertical space, instead of a
    # fixed 1.6in row sitting near the top with a big blank stretch
    # below it — a 2-step process used to look like a stub.
    content_top_in, content_bottom_in = 1.7, 7.0
    chip_h_in = 1.6 if n >= 5 else (2.7 if n <= 2 else 2.1)
    chip_h = Inches(chip_h_in)
    area_top = Inches(content_top_in + 0.3 + max(0, (content_bottom_in - content_top_in - 0.3 - chip_h_in) / 2))
    chip_w = (area_w - arrow_w * (n - 1)) / n

    x = area_left
    for i, step in enumerate(steps):
        text = step.get("text", "") if isinstance(step, dict) else str(step)

        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, area_top, chip_w, chip_h)
        chip.adjustments[0] = 0.15
        _set_fill(chip, accent)
        chip.shadow.inherit = False

        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + chip_w / 2 - Inches(0.22), area_top - Inches(0.22), Inches(0.44), Inches(0.44))
        _set_fill(num, WHITE)
        num.line.color.rgb = NAVY
        num.line.width = Pt(1.5)
        num.shadow.inherit = False
        num_tf = num.text_frame
        num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        num_tf.text = str(i + 1)
        num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_run = num_tf.paragraphs[0].runs[0]
        num_run.font.size = Pt(14)
        num_run.font.bold = True
        num_run.font.color.rgb = NAVY

        text_box = slide.shapes.add_textbox(x + Inches(0.1), area_top + Inches(0.25), chip_w - Inches(0.2), chip_h - Inches(0.35))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = text
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE if accent != GOLD_LIGHT else NAVY

        x += chip_w
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, area_top + chip_h / 2 - Inches(0.14), arrow_w, Inches(0.28))
            _set_fill(arrow, NAVY_MID)
            arrow.shadow.inherit = False
            x += arrow_w


def _add_comparison_slide(prs: Presentation, layout, heading: str, left: dict, right: dict, index: int) -> None:
    """Pros/cons, before/after, X vs Y — two colored columns side by
    side instead of one paragraph trying to hold both sides at once."""
    slide = prs.slides.add_slide(layout)
    _add_slide_chrome(prs, slide, heading, index, GOLD)

    col_top = Inches(1.85)
    col_h = Inches(5.15)
    gap = Inches(0.35)
    area_left = Inches(1.55)
    area_w = prs.slide_width - area_left - Inches(0.6)
    col_w = (area_w - gap) / 2

    panels = [(left, NAVY, area_left), (right, GOLD, area_left + col_w + gap)]
    for panel, color, x in panels:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_top, col_w, col_h)
        card.adjustments[0] = 0.05
        _set_fill(card, color)
        card.shadow.inherit = False

        label_box = slide.shapes.add_textbox(x + Inches(0.3), col_top + Inches(0.22), col_w - Inches(0.6), Inches(0.5))
        label_tf = label_box.text_frame
        label_tf.text = panel.get("label", "")
        label_run = label_tf.paragraphs[0].runs[0]
        label_run.font.size = Pt(18)
        label_run.font.bold = True
        label_run.font.color.rgb = WHITE if color != GOLD else NAVY

        # Item spacing scales to spread the list across the FULL column
        # height (see _distribute_fill) — a 2-3 item side used to leave
        # a big blank stretch under the last item while the panel's
        # colored background kept going.
        items_list = panel.get("items", [])
        label_h_in = 0.85
        avail_h_in = (col_h / 914400) - label_h_in - 0.3  # EMU -> inches, minus bottom padding
        item_h_in, offset_in = _distribute_fill(len(items_list), avail_h_in, 0.0, min_item=0.5, max_item=1.0)
        item_top = col_top + Inches(label_h_in + offset_in)
        for item in items_list:
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), item_top + Inches(0.08), Inches(0.12), Inches(0.12))
            marker.fill.solid()
            marker.fill.fore_color.rgb = WHITE if color != GOLD else NAVY
            marker.line.fill.background()
            marker.shadow.inherit = False

            item_box = slide.shapes.add_textbox(x + Inches(0.58), item_top - Inches(0.08), col_w - Inches(0.9), Inches(item_h_in))
            item_tf = item_box.text_frame
            item_tf.word_wrap = True
            item_tf.vertical_anchor = MSO_ANCHOR.TOP
            item_tf.text = item
            item_run = item_tf.paragraphs[0].runs[0]
            item_run.font.size = Pt(14)
            item_run.font.color.rgb = WHITE if color != GOLD else NAVY
            item_run.font.name = "Arial"
            item_top += Inches(item_h_in)


def _set_transparency(shape, alpha_pct: int) -> None:
    """python-pptx has no first-class alpha API — this pokes the
    <a:alpha> element directly into the shape's solid fill XML.
    alpha_pct is how OPAQUE the shape should look (0 = invisible, 100 = solid)."""
    sp = shape.fill.fore_color._xFill
    alpha = sp.find(qn("a:srgbClr"))
    if alpha is None:
        return
    a = alpha.makeelement(qn("a:alpha"), {"val": str(alpha_pct * 1000)})
    alpha.append(a)
